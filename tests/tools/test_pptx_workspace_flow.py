from __future__ import annotations

from pathlib import Path
import zipfile

import pytest

from deeptutor.services.path_service import PathService
from deeptutor.services.sandbox.backends import RestrictedSubprocessBackend
from deeptutor.services.session.artifact_attachments import fill_preview_text
from deeptutor.services.workspace import ContentWorkspaceService
from deeptutor.services.workspace.execution import prepare_workspace_execution_env
from deeptutor.tools.exec_tool import ExecTool
from deeptutor.tools.workspace import WorkspacePresentTool


def test_pptx_skill_teaches_only_unified_exec_and_logical_output_paths() -> None:
    root = Path(__file__).parents[2]
    skill = (root / "deeptutor" / "skills" / "builtin" / "pptx" / "SKILL.md").read_text()

    for removed_tool in ("code_execution", "run_code", "code_execute"):
        assert removed_tool not in skill
    assert "Use `exec` with complete Python source" in skill
    assert "`language: python`" in skill
    assert "**User workspace**" in skill
    assert "DEEPTUTOR_WORKSPACE_ROOT" not in skill
    assert "python -c" not in skill
    assert "heredoc" not in skill
    assert "$PWD" not in skill
    assert "command -v" not in skill
    assert "with your image-view capability" not in skill
    assert "`workspace_present` only presents a file to the user" in skill


def test_pptx_dependency_is_packaged_for_pip_source_and_docker_runner() -> None:
    root = Path(__file__).parents[2]

    assert '"python-pptx>=1.0.0"' in (root / "pyproject.toml").read_text()
    assert "python-pptx" in (root / "Dockerfile.runner").read_text()


@pytest.mark.asyncio
async def test_real_pptx_exec_stays_in_outputs_and_is_presented(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    """Exercise the same Python runtime, artifact collector, and presentation path as chat."""

    from deeptutor.services.workspace import service as workspace_service_module

    runtime_paths = PathService(workspace_root=tmp_path / "runtime")
    runtime_paths.ensure_all_directories()
    monkeypatch.setattr(workspace_service_module, "get_path_service", lambda: runtime_paths)

    workspace_root = tmp_path / "chosen-workspace"
    workspace_root.mkdir()
    monkeypatch.setenv("DEEPTUTOR_WORKSPACE_ROOT", str(workspace_root))
    monkeypatch.setenv("DEEPTUTOR_WORKSPACE_ALLOWED_ROOTS", str(tmp_path))

    workspace_service = ContentWorkspaceService()
    binding = workspace_service.current_binding(ensure_output=True)
    turn_dir = workspace_root / "outputs" / "chat" / "session" / "turn"
    code_workdir = turn_dir / "code_runs"
    code_workdir.mkdir(parents=True)

    class DirectSandboxService:
        async def run(self, request, *, user_id: str):
            return await RestrictedSubprocessBackend().exec(request)

    import deeptutor.services.sandbox as sandbox_package
    import deeptutor.services.workspace as workspace_package
    import deeptutor.tools.workspace as workspace_tools

    monkeypatch.setattr(sandbox_package, "get_sandbox_service", lambda: DirectSandboxService())
    monkeypatch.setattr(
        workspace_package,
        "get_content_workspace_service",
        lambda: workspace_service,
    )
    monkeypatch.setattr(
        workspace_tools,
        "get_content_workspace_service",
        lambda: workspace_service,
    )

    first_result = await ExecTool().execute(
        code=(
            "from pptx import Presentation\n"
            "from pptx.util import Inches\n"
            "import zipfile\n"
            "prs = Presentation()\n"
            "slide = prs.slides.add_slide(prs.slide_layouts[0])\n"
            "slide.shapes.title.text = 'DeepTutor PPTX smoke test'\n"
            "slide.placeholders[1].text = 'Generated through unified exec'\n"
            "path = 'deeptutor_smoke.pptx'\n"
            "prs.save(path)\n"
            "reopened = Presentation(path)\n"
            "assert len(reopened.slides) == 1\n"
            "assert reopened.slides[0].shapes.title.text == 'DeepTutor PPTX smoke test'\n"
            "with zipfile.ZipFile(path) as package:\n"
            "    assert 'ppt/presentation.xml' in package.namelist()\n"
            "print('verified deeptutor_smoke.pptx')\n"
        ),
        language="python",
        _sandbox_user_id="pptx-test",
        _sandbox_code_workdir=str(code_workdir),
        _sandbox_env=prepare_workspace_execution_env(turn_dir, workspace_root=workspace_root),
        _workspace_id=binding.workspace_id,
        _workspace_root=str(workspace_root),
    )
    first_artifact = first_result.metadata["artifacts"][0]
    assert first_result.metadata["workspace_items"] == []
    assert first_result.sources == []
    first_presentation = await WorkspacePresentTool().execute(
        items=[{"path": first_artifact["relative_path"]}],
        _workspace_id=binding.workspace_id,
    )

    # A second source call uses the same turn-local working directory, so the
    # model can revise the deck without learning or reconstructing a host path.
    result = await ExecTool().execute(
        code=(
            "from pptx import Presentation\n"
            "prs = Presentation('deeptutor_smoke.pptx')\n"
            "slide = prs.slides.add_slide(prs.slide_layouts[5])\n"
            "slide.shapes.title.text = 'Verified revision'\n"
            "prs.save('deeptutor_smoke.pptx')\n"
            "reopened = Presentation('deeptutor_smoke.pptx')\n"
            "assert len(reopened.slides) == 2\n"
            "print('revised deeptutor_smoke.pptx')\n"
        ),
        language="python",
        _sandbox_user_id="pptx-test",
        _sandbox_code_workdir=str(code_workdir),
        _sandbox_env=prepare_workspace_execution_env(turn_dir, workspace_root=workspace_root),
        _workspace_id=binding.workspace_id,
        _workspace_root=str(workspace_root),
    )

    deck = code_workdir / "deeptutor_smoke.pptx"
    assert deck.is_file()
    deck.relative_to(workspace_root / "outputs")
    assert not (workspace_root / "deeptutor_smoke.pptx").exists()
    assert result.success is True
    assert deck.stat().st_size > 1_000
    assert zipfile.is_zipfile(deck)
    with zipfile.ZipFile(deck) as package:
        assert package.testzip() is None
        assert "[Content_Types].xml" in package.namelist()
        assert "ppt/slides/slide1.xml" in package.namelist()
        assert "ppt/slides/slide2.xml" in package.namelist()

    artifacts = result.metadata["artifacts"]
    assert len(artifacts) == 1
    assert artifacts[0]["filename"] == "deeptutor_smoke.pptx"
    assert artifacts[0]["relative_path"].startswith("outputs/")
    assert artifacts[0]["mime_type"] == (
        "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )

    assert result.metadata["workspace_items"] == []
    assert result.sources == []
    assert "not yet presented" in result.content
    presentation = await WorkspacePresentTool().execute(
        items=[{"path": artifacts[0]["relative_path"]}],
        _workspace_id=binding.workspace_id,
    )
    presented = presentation.metadata["workspace_items"]
    assert len(presented) == 1
    assert presented[0]["relative_path"] == artifacts[0]["relative_path"]
    assert presented[0]["generated"] is True
    assert presented[0]["url"].startswith("/files/workspace-items/")
    assert presentation.sources == [{"type": "workspace_item", **presented[0]}]
    assert str(workspace_root) not in result.content
    assert (
        first_presentation.metadata["workspace_items"][0]["workspace_item_id"]
        != presented[0]["workspace_item_id"]
    )

    preview_attachments = [
        {
            "filename": presented[0]["filename"],
            "url": presented[0]["url"],
            "mime_type": presented[0]["mime_type"],
        }
    ]
    await fill_preview_text(preview_attachments)
    assert "DeepTutor PPTX smoke test" in preview_attachments[0]["extracted_text"]
    assert "Verified revision" in preview_attachments[0]["extracted_text"]
